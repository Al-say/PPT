from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from bs4 import BeautifulSoup
import os
import requests
from io import BytesIO
import re

def add_decorative_shape(slide, shape_type=MSO_SHAPE.RECTANGLE, left=0, top=0, width=1, height=1):
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(200, 0, 0)  # Red color for Chinese theme
    shape.line.color.rgb = RGBColor(150, 0, 0)
    return shape

def create_title_slide(prs, title, subtitle=""):
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    
    # Add decorative elements
    add_decorative_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.5, 0.5, 1, 15)
    add_decorative_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 14.5, 0.5, 1, 15)
    
    # Find title and subtitle placeholders
    title_placeholder = None
    subtitle_placeholder = None
    
    for shape in slide.placeholders:
        if shape.placeholder_format.type == 1:  # Title
            title_placeholder = shape
        elif shape.placeholder_format.type == 2:  # Subtitle
            subtitle_placeholder = shape
    
    # Add and format title
    if title_placeholder:
        title_placeholder.text = title
        title_frame = title_placeholder.text_frame
        paragraph = title_frame.paragraphs[0]
        paragraph.font.size = Pt(44)
        paragraph.font.name = "Microsoft YaHei"
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(150, 0, 0)
    
    # Add and format subtitle if provided
    if subtitle and subtitle_placeholder:
        subtitle_placeholder.text = subtitle
        subtitle_frame = subtitle_placeholder.text_frame
        paragraph = subtitle_frame.paragraphs[0]
        paragraph.font.size = Pt(32)
        paragraph.font.name = "Microsoft YaHei"
        paragraph.font.color.rgb = RGBColor(100, 0, 0)
    
    return slide

def create_content_slide(prs, title, content, image_url=None, has_bullet_points=False, include_decoration=True):
    # Choose layout based on content type
    layout_idx = 1 if not has_bullet_points else 2
    content_slide_layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(content_slide_layout)
    
    # Add decorative elements if requested
    if include_decoration:
        add_decorative_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.5, 0.5, 0.2, 8)
    
    # Add and format title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_frame = title_shape.text_frame
    paragraph = title_frame.paragraphs[0]
    paragraph.font.size = Pt(36)
    paragraph.font.name = "Microsoft YaHei"
    paragraph.font.bold = True
    paragraph.font.color.rgb = RGBColor(150, 0, 0)
    
    # Add and format content
    content_shape = slide.shapes.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.clear()
    
    if has_bullet_points and isinstance(content, (list, tuple)):
        for point in content:
            p = text_frame.add_paragraph()
            p.text = "• " + point
            p.font.size = Pt(24)
            p.font.name = "Microsoft YaHei"
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.space_after = Pt(12)  # Add some spacing between points
    else:
        p = text_frame.add_paragraph()
        p.text = content
        p.font.size = Pt(24)
        p.font.name = "Microsoft YaHei"
        p.font.color.rgb = RGBColor(0, 0, 0)
    
    # Add image if URL is provided
    if image_url:
        try:
            # Download image
            response = requests.get(image_url)
            image_stream = BytesIO(response.content)
            
            # Add image to slide
            slide.shapes.add_picture(
                image_stream,
                Inches(1),  # Left position
                Inches(3),  # Top position
                width=Inches(4)  # Width, height will maintain aspect ratio
            )
        except Exception as e:
            print(f"Failed to add image {image_url}: {str(e)}")
    
    return slide

def extract_image_url(soup):
    """Extract image URL from the HTML content"""
    # Look for img tags
    img_tag = soup.find('img')
    if img_tag and 'src' in img_tag.attrs:
        src = img_tag['src']
        # Handle both absolute and relative URLs
        if src.startswith(('http://', 'https://')):
            return src
        elif src.startswith('//'):
            return f'https:{src}'
    
    # Look for background-image in style attributes
    elements_with_style = soup.find_all(lambda tag: tag.get('style', '') and 'background-image' in tag['style'])
    for element in elements_with_style:
        style = element['style']
        url_match = re.search(r'background-image:\s*url\([\'"]?([^\'"]*)[\'"]?\)', style)
        if url_match:
            url = url_match.group(1)
            if url.startswith(('http://', 'https://')):
                return url
            elif url.startswith('//'):
                return f'https:{url}'
    
    return None

def extract_bullet_points(soup):
    """Extract bullet points from HTML content if available"""
    bullet_points = []
    ul = soup.find('ul')
    if ul:
        for li in ul.find_all('li'):
            bullet_points.append(li.text.strip())
    return bullet_points

def main():
    # Initialize presentation with widescreen format and theme colors
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # Set theme colors (red theme for Chinese presentation)
    theme_colors = [
        (RGBColor(150, 0, 0), "Accent 1"),
        (RGBColor(200, 0, 0), "Accent 2"),
        (RGBColor(100, 0, 0), "Accent 3"),
        (RGBColor(50, 0, 0), "Accent 4"),
    ]
    
    # Create title slide with background
    create_title_slide(prs, "五四运动：青年觉醒的时代号角", "从爱国运动到青年节，探寻中国青年的精神传承")
    
    # Process sections
    sections_dir = "sections"
    for filename in sorted(os.listdir(sections_dir)):
        if filename.endswith(".html"):
            with open(os.path.join(sections_dir, filename), 'r', encoding='utf-8') as f:
                content = f.read()
                soup = BeautifulSoup(content, 'html.parser')
                
                # Extract heading and paragraphs
                heading = soup.find('h1')
                paragraphs = soup.find_all('p')
                
                if heading:
                    # Extract content elements
                    title = heading.text
                    bullet_points = extract_bullet_points(soup)
                    
                    if bullet_points:
                        # Create bullet point slide
                        create_content_slide(prs, title, bullet_points, has_bullet_points=True)
                    else:
                        # Create regular content slide with paragraphs
                        content = '\n'.join([p.text for p in paragraphs if p.text])
                        if content:
                            # Extract image URL if available
                            image_url = extract_image_url(soup)
                            create_content_slide(prs, title, content, image_url)
                        else:
                            # For section headers without content
                            create_title_slide(prs, title)
    
    # Save the presentation
    prs.save('五四运动.pptx')

if __name__ == "__main__":
    main()
